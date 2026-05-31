import streamlit as st

import pymupdf4llm

import mammoth

import tempfile

import os

import re

import zipfile

import io

from pathlib import Path

from datetime import datetime

from pptx import Presentation

# ── Page config ───────────────────────────────────────────────────────────────

st.set_page_config(

page_title="Docs → Markdown",

page_icon="📄",

layout="centered",

)

# ── Styles ────────────────────────────────────────────────────────────────────

st.markdown("""

<style>

@import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;600&family=IBM+Plex+Sans:wght@300;400;600&display=swap');

html, body, [class*="css"] { font-family: 'IBM Plex Sans', sans-serif; }

.stApp { background: #0d0d0d; color: #e8e8e8; }

header[data-testid="stHeader"] { background: transparent; }

h1 {

font-family: 'IBM Plex Mono', monospace !important;

font-size: 1.8rem !important;

color: #f0f0f0 !important;

letter-spacing: -0.02em;

margin-bottom: 0 !important;

}

.subtitle {

font-family: 'IBM Plex Mono', monospace;

font-size: 0.75rem;

color: #555;

margin-bottom: 2rem;

letter-spacing: 0.08em;

}

[data-testid="stFileUploader"] {

background: #141414;

border: 1px solid #2a2a2a;

border-radius: 6px;

padding: 0.5rem;

}

[data-testid="stFileUploader"]:hover { border-color: #444; }

[data-testid="stMetric"] {

background: #141414;

border: 1px solid #222;

border-radius: 6px;

padding: 0.75rem 1rem;

}

[data-testid="stMetricLabel"] {

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.65rem !important;

color: #555 !important;

text-transform: uppercase;

letter-spacing: 0.1em;

}

[data-testid="stMetricValue"] {

font-family: 'IBM Plex Mono', monospace !important;

font-size: 1.1rem !important;

color: #c8f09a !important;

}

[data-testid="stMetricDelta"] {

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.7rem !important;

}

/* Primary button */

div[data-testid="stButton"] > button[kind="primary"] {

background: #c8f09a !important;

color: #0d0d0d !important;

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.85rem !important;

font-weight: 600 !important;

letter-spacing: 0.05em;

border: none !important;

border-radius: 4px !important;

padding: 0.6rem 2rem !important;

width: 100%;

}

div[data-testid="stButton"] > button[kind="primary"]:hover { background: #d8f5b0 !important; }

/* Secondary button */

div[data-testid="stButton"] > button[kind="secondary"] {

background: transparent !important;

color: #555 !important;

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.8rem !important;

border: 1px solid #2a2a2a !important;

border-radius: 4px !important;

padding: 0.6rem 2rem !important;

width: 100%;

}

div[data-testid="stButton"] > button[kind="secondary"]:hover {

border-color: #555 !important;

color: #ccc !important;

}

/* Download button */

[data-testid="stDownloadButton"] > button {

background: #c8f09a !important;

color: #0d0d0d !important;

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.8rem !important;

font-weight: 600 !important;

letter-spacing: 0.05em;

border: none !important;

border-radius: 4px !important;

padding: 0.5rem 1.5rem !important;

width: 100%;

}

[data-testid="stDownloadButton"] > button:hover { background: #d8f5b0 !important; }

.file-row {

background: #111;

border: 1px solid #1e1e1e;

border-radius: 6px;

padding: 0.75rem 1rem;

margin-bottom: 0.5rem;

display: flex;

align-items: center;

justify-content: space-between;

font-family: 'IBM Plex Mono', monospace;

font-size: 0.75rem;

}

.file-name { color: #ccc; }

.file-ok { color: #c8f09a; }

.file-err { color: #f09a9a; }

.file-kb { color: #555; margin-left: 0.5rem; }

/* History table */

.hist-row {

background: #0e0e0e;

border: 1px solid #1a1a1a;

border-radius: 5px;

padding: 0.6rem 1rem;

margin-bottom: 0.35rem;

display: grid;

grid-template-columns: 2fr 0.7fr 0.7fr 0.7fr 0.9fr;

align-items: center;

font-family: 'IBM Plex Mono', monospace;

font-size: 0.7rem;

gap: 0.5rem;

}

.hist-header {

color: #444;

text-transform: uppercase;

letter-spacing: 0.08em;

font-size: 0.62rem;

border-bottom: 1px solid #1e1e1e;

padding-bottom: 0.4rem;

margin-bottom: 0.4rem;

}

.hist-name { color: #aaa; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }

.hist-ext { color: #555; }

.hist-kb { color: #555; }

.hist-tok { color: #c8f09a; }

.hist-time { color: #444; text-align: right; }

.badge-pdf { color: #f0b49a; }

.badge-docx { color: #9ab8f0; }

.badge-pptx { color: #c49af0; }

[data-testid="stExpander"] {

background: #111 !important;

border: 1px solid #1e1e1e !important;

border-radius: 6px !important;

}

/* text_area editor */

.stTextArea textarea {

background: #111 !important;

border: 1px solid #222 !important;

border-radius: 6px !important;

font-family: 'IBM Plex Mono', monospace !important;

font-size: 0.72rem !important;

color: #888 !important;

line-height: 1.7 !important;

}

.stTextArea textarea:focus {

border-color: #444 !important;

color: #ccc !important;

}

hr { border-color: #1e1e1e !important; }

.stAlert { background: #111 !important; border-radius: 6px !important; }

</style>

""", unsafe_allow_html=True)

# ── Session state ─────────────────────────────────────────────────────────────

defaults = {

"uploader_key": 0,

"results": None,

"history": [], # list of dicts — one per successful conversion

"edited_md": {}, # md_name → edited text

}

for k, v in defaults.items():

if k not in st.session_state:

st.session_state[k] = v

# ── Converters ────────────────────────────────────────────────────────────────

def estimate_tokens(text: str) -> int:

return len(text) // 4

def pdf_bytes_to_md(pdf_bytes: bytes) -> str:

with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:

tmp.write(pdf_bytes)

tmp_path = tmp.name

try:

md = pymupdf4llm.to_markdown(tmp_path)

finally:

os.unlink(tmp_path)

return md

def docx_bytes_to_md(docx_bytes: bytes) -> str:

buf = io.BytesIO(docx_bytes)

style_map = """

p[style-name='Heading 1'] => h1:fresh

p[style-name='Heading 2'] => h2:fresh

p[style-name='Heading 3'] => h3:fresh

p[style-name='Heading 4'] => h4:fresh

p[style-name='Heading 5'] => h5:fresh

p[style-name='Heading 6'] => h6:fresh

r[style-name='Strong'] => strong

r[style-name='Emphasis'] => em

"""

result = mammoth.convert_to_markdown(buf, style_map=style_map)

md = re.sub(r'\n{3,}', '\n\n', result.value).strip()

return md

# FIX #1: @st.cache_data evita re-procesar el PPTX en cada re-ejecución de Streamlit.
# FIX #2: ph_fmt guarda la referencia a placeholder_format para verificar None
#          antes de acceder a .idx, eliminando el error "shape is not a placeholder".
@st.cache_data(show_spinner=False)
def pptx_bytes_to_md(pptx_bytes: bytes) -> str:

    buf = io.BytesIO(pptx_bytes)

    prs = Presentation(buf)

    slides_md = []

    for slide_num, slide in enumerate(prs.slides, start=1):

        lines = []

        title_text = ""

        body_parts = []

        table_parts = []

        notes_text = ""

        for shape in slide.shapes:

            if shape.has_table:

                table = shape.table

                rows_md = []

                for r_idx, row in enumerate(table.rows):

                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]

                    rows_md.append("| " + " | ".join(cells) + " |")

                    if r_idx == 0:

                        rows_md.append("|" + "|".join(["---"] * len(cells)) + "|")

                table_parts.append("\n".join(rows_md))

                continue

            if not shape.has_text_frame:

                continue

            # FIX #2: guardar referencia antes de acceder a .idx
            ph_fmt = shape.placeholder_format
            is_title = (
                ph_fmt is not None
                and ph_fmt.idx in (0, 1)
            )

            frame_lines = []

            for para in shape.text_frame.paragraphs:

                para_text = para.text.strip()

                if not para_text:

                    continue

                level = para.level if para.level is not None else 0

                indent = " " * level

                all_bold = (

                    all(run.font.bold is True for run in para.runs if run.text.strip())

                    and para.runs

                )

                if all_bold and level == 0:

                    frame_lines.append(f"**{para_text}**")

                else:

                    frame_lines.append(f"{indent}- {para_text}")

            if is_title:

                title_text = shape.text_frame.text.strip()

            else:

                body_parts.extend(frame_lines)

        if slide.has_notes_slide:

            notes_tf = slide.notes_slide.notes_text_frame

            raw_notes = notes_tf.text.strip() if notes_tf else ""

            if raw_notes and raw_notes != str(slide_num):

                notes_text = raw_notes

        heading = f"## Slide {slide_num}" + (f": {title_text}" if title_text else "")

        lines.append(heading)

        if body_parts:

            lines.append("")

            lines.extend(body_parts)

        if table_parts:

            lines.append("")

            for t in table_parts:

                lines.append(t)

        if notes_text:

            lines.append("")

            lines.append(f"> **Notes:** {notes_text}")

        slides_md.append("\n".join(lines))

    return "\n\n---\n\n".join(slides_md)

FILE_ICONS = {

"pdf": ("📄", "badge-pdf"),

"docx": ("📝", "badge-docx"),

"pptx": ("📊", "badge-pptx"),

}

def convert_file(file_bytes: bytes, extension: str) -> str:

if extension == "pdf":

return pdf_bytes_to_md(file_bytes)

elif extension == "docx":

return docx_bytes_to_md(file_bytes)

elif extension == "pptx":

return pptx_bytes_to_md(file_bytes)

else:

raise ValueError(f"Unsupported format: {extension}")

def build_zip(results: list) -> bytes:

buf = io.BytesIO()

with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:

for name, md_text in results:

zf.writestr(name, md_text.encode("utf-8"))

return buf.getvalue()

def do_clear():

st.session_state.uploader_key += 1

st.session_state.results = None

st.session_state.edited_md = {}

# ── Header ────────────────────────────────────────────────────────────────────

st.markdown("# docs → markdown")

st.markdown(

'<p class="subtitle">TOKEN-EFFICIENT CONVERSION FOR LLM CONSUMPTION · PDF · DOCX · PPTX</p>',

unsafe_allow_html=True,

)

st.divider()

# ── Upload ────────────────────────────────────────────────────────────────────

uploaded_files = st.file_uploader(

"Upload PDFs, Word documents, or PowerPoint presentations",

type=["pdf", "docx", "pptx"],

accept_multiple_files=True,

label_visibility="collapsed",

key=f"uploader_{st.session_state.uploader_key}",

)

st.markdown("")

col_convert, col_clear = st.columns([3, 1])

with col_convert:

convert_clicked = st.button("⚡ Convert", type="primary", disabled=not uploaded_files)

with col_clear:

st.button(

"🗑 Clear", type="secondary", on_click=do_clear,

disabled=not uploaded_files and st.session_state.results is None,

)

# ── Conversion ────────────────────────────────────────────────────────────────

if convert_clicked and uploaded_files:

results = []

total_file_kb = 0

total_md_kb = 0

total_tokens = 0

progress = st.progress(0, text="Converting files…")

for i, f in enumerate(uploaded_files):

file_bytes = f.read()

file_kb = len(file_bytes) / 1024

total_file_kb += file_kb

ext = Path(f.name).suffix.lstrip(".").lower()

try:

md_text = convert_file(file_bytes, ext)

md_kb = len(md_text.encode("utf-8")) / 1024

tokens = estimate_tokens(md_text)

total_md_kb += md_kb

total_tokens += tokens

stem = Path(f.name).stem

entry = {

"md_name": f"{stem}.md",

"md_text": md_text,

"file_kb": file_kb,

"md_kb": md_kb,

"tokens": tokens,

"ok": True,

"orig_name": f.name,

"ext": ext,

"converted_at": datetime.now().strftime("%H:%M:%S"),

}

results.append(entry)

# ── Add to history (avoid duplicates by md_name) ──────────────────

existing_names = [h["md_name"] for h in st.session_state.history]

if entry["md_name"] not in existing_names:

st.session_state.history.insert(0, {

"md_name": entry["md_name"],

"orig_name": entry["orig_name"],

"ext": ext,

"file_kb": file_kb,

"md_kb": md_kb,

"tokens": tokens,

"converted_at": entry["converted_at"],

})

# Initialise editor state with original text

st.session_state.edited_md.setdefault(entry["md_name"], md_text)

except Exception as e:

results.append({

"md_name": f.name,

"ok": False,

"error": str(e),

"orig_name": f.name,

"file_kb": file_kb,

"ext": ext,

"converted_at": datetime.now().strftime("%H:%M:%S"),

})

progress.progress((i + 1) / len(uploaded_files),

text=f"Converting {i+1}/{len(uploaded_files)}…")

progress.empty()

st.session_state.results = {

"files": results,

"total_file_kb": total_file_kb,

"total_md_kb": total_md_kb,

"total_tokens": total_tokens,

}

# ── Results ───────────────────────────────────────────────────────────────────

if st.session_state.results:

data = st.session_state.results

results = data["files"]

ok_list = [r for r in results if r["ok"]]

st.divider()

# Status rows

for r in results:

icon, badge_class = FILE_ICONS.get(r["ext"], ("📄", "file-ok"))

if r["ok"]:

st.markdown(

f'<div class="file-row">'

f'<span class="file-name"><span class="{badge_class}">{icon}</span> {r["orig_name"]}</span>'

f'<span><span class="file-ok">✓ converted</span>'

f'<span class="file-kb">{r["file_kb"]:.1f} KB → {r["md_kb"]:.1f} KB'

f' · ~{r["tokens"]:,} tokens</span></span>'

f'</div>',

unsafe_allow_html=True,

)

else:

st.markdown(

f'<div class="file-row">'

f'<span class="file-name"><span class="{badge_class}">{icon}</span> {r["orig_name"]}</span>'

f'<span class="file-err">✗ {r["error"]}</span>'

f'</div>',

unsafe_allow_html=True,

)

# Summary metrics

if ok_list:

reduction = (1 - data["total_md_kb"] / data["total_file_kb"]) * 100

st.divider()

c1, c2, c3, c4 = st.columns(4)

with c1: st.metric("Files", f"{len(ok_list)}/{len(results)}")

with c2: st.metric("Total size", f"{data['total_file_kb']:.1f} KB")

with c3: st.metric("Reduction", f"{reduction:.0f}%",

delta=f"-{data['total_file_kb'] - data['total_md_kb']:.1f} KB")

with c4: st.metric("~Total tokens", f"{data['total_tokens']:,}")

st.divider()

# ── Download + editor per file ────────────────────────────────────────

for r in ok_list:

md_name = r["md_name"]

icon, _ = FILE_ICONS.get(r["ext"], ("📄", ""))

with st.expander(f"{icon} {md_name}", expanded=(len(ok_list) == 1)):

# ── Inline editor ────────────────────────────────────────────

editor_key = f"editor_{md_name}"

if editor_key not in st.session_state:

st.session_state[editor_key] = r["md_text"]

st.text_area(

"Edit Markdown before downloading",

height=320,

key=editor_key,

label_visibility="collapsed",

)

# ── Download (uses edited text) ──────────────────────────────

st.download_button(

label="⬇ Download .md",

data=st.session_state[editor_key].encode("utf-8"),

file_name=md_name,

mime="text/markdown",

key=f"dl_{md_name}",

)

# Bulk download always uses edited versions

if len(ok_list) > 1:

st.divider()

edited_pairs = [

(r["md_name"], st.session_state.get(f"editor_{r['md_name']}", r["md_text"]))

for r in ok_list

]

zip_bytes = build_zip(edited_pairs)

st.download_button(

label=f"⬇ Download all ({len(ok_list)} files) as .zip",

data=zip_bytes,

file_name="markdown_files.zip",

mime="application/zip",

key="dl_zip",

)

elif not uploaded_files:

st.markdown("")

st.info("Drag and drop PDFs, Word docs, or PowerPoint files above, or click to browse.", icon="📎")

# ── Session history ───────────────────────────────────────────────────────────

if st.session_state.history:

st.divider()

with st.expander(f"🕓 Session history ({len(st.session_state.history)} files)", expanded=False):

# Header row

st.markdown(

'<div class="hist-row hist-header">'

'<span>File</span><span>Type</span><span>Size</span>'

'<span>~Tokens</span><span style="text-align:right">Time</span>'

'</div>',

unsafe_allow_html=True,

)

for h in st.session_state.history:

icon, badge = FILE_ICONS.get(h["ext"], ("📄", ""))

st.markdown(

f'<div class="hist-row">'

f'<span class="hist-name" title="{h["orig_name"]}">{h["orig_name"]}</span>'

f'<span class="hist-ext"><span class="{badge}">{icon} {h["ext"].upper()}</span></span>'

f'<span class="hist-kb">{h["file_kb"]:.1f}→{h["md_kb"]:.1f} KB</span>'

f'<span class="hist-tok">~{h["tokens"]:,}</span>'

f'<span class="hist-time">{h["converted_at"]}</span>'

f'</div>',

unsafe_allow_html=True,

)

# Aggregate stats

total_tok = sum(h["tokens"] for h in st.session_state.history)

total_in = sum(h["file_kb"] for h in st.session_state.history)

total_out = sum(h["md_kb"] for h in st.session_state.history)

st.caption(

f"Session total · {len(st.session_state.history)} files · "

f"{total_in:.1f} KB in → {total_out:.1f} KB out · ~{total_tok:,} tokens"

)

if st.button("🗑 Clear history", type="secondary", key="clear_hist"):

st.session_state.history = []

st.rerun()
